"""
异步PPT翻译模块
基于ppt_translate.py的逻辑，支持并行处理幻灯片和异步翻译，提高处理效率
保持与原始功能的完全兼容性，包括表格处理、样式保持等
"""
import os
import sys
import time
import asyncio
import logging
import re
import json
import platform
from typing import Dict, List, Any, Optional, Union, Tuple
import concurrent.futures
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
import difflib

# 导入异步API客户端
from .local_qwen_async import translate_async, batch_translate_async, get_field_async
from ..utils.thread_pool_executor import thread_pool, TaskType
from ..utils.enhanced_task_queue import translation_queue

# 导入基于页面的翻译机制
from .page_based_translation import translate_slide_by_page, get_translation_statistics

# 导入复杂形状处理函数和内容检测函数
from .ppt_translate import (
    detect_complex_shape_type,
    save_complex_shape_properties,
    restore_complex_shape_properties,
    has_shape_deformed,
    safe_set_autofit_with_size_preservation,
    has_meaningful_text_content,
    should_adjust_textbox_layout,
    get_textbox_content_summary,
    safe_set_autofit_with_content_check
)

# 配置日志记录器
logger = logging.getLogger(__name__)

# 优化: 分块大小，用于批量处理
BATCH_SIZE = 20  # 每批处理的文本数量
MAX_BATCH_CHAR_COUNT = 2000  # 每批最大字符数
SLIDE_PROCESSING_THREADS = 3  # 幻灯片并行处理线程数（减少以避免资源竞争）

def get_font_color(run):
    """获取文本颜色，处理RGB颜色和主题颜色"""
    color_format = run.font.color
    if color_format.type == 1:  # MSO_COLOR_TYPE.RGB (1)
        return color_format.rgb
    elif color_format.type == 2:  # MSO_COLOR_TYPE.SCHEME (2)
        return color_format.theme_color
    return None

def apply_font_color(run, color):
    """应用文本颜色，处理RGB颜色和主题颜色"""
    if isinstance(color, RGBColor):
        run.font.color.rgb = color
    else:
        run.font.color.theme_color = color

def compare_strings_ignore_spaces(str1, str2):
    """比较两个字符串，忽略空格"""
    str1_no_spaces = str1.replace(" ", "")
    str2_no_spaces = str2.replace(" ", "")
    return str1_no_spaces == str2_no_spaces

def find_most_similar(target, candidates):
    """从候选列表中找到与目标最相似的字符串"""
    if not candidates:
        return None
    return max(candidates, key=lambda x: difflib.SequenceMatcher(None, target, x).ratio())

def remove_invalid_utf8_chars(s: str) -> str:
    """移除字符串中无效的UTF-8字符"""
    utf8_bytes = s.encode('utf-8', errors='ignore')
    clean_str = utf8_bytes.decode('utf-8', errors='ignore')
    return clean_str

def is_valid_reference(text):
    """检查文本是否为有效的参考文献"""
    pattern = r'\d+\s*[A-Za-z&\s\.\-]+,\s*\d{4}'
    return bool(re.match(pattern, text))

def is_page_number(text):
    """检查文本是否为页码"""
    text = text.strip()
    if re.fullmatch(r'\d{1,3}', text):
        return True
    return False

async def _adjust_ppt_layout_async(presentation_path: str) -> bool:
    """
    异步调整PPT布局，使用现有的set_textbox_autofit函数

    Args:
        presentation_path: PPT文件路径

    Returns:
        调整是否成功
    """
    try:
        # 检查操作系统，只在Windows上执行COM操作
        if platform.system() != "Windows":
            logger.info("非Windows系统，跳过COM布局调整，使用基础布局调整")
            return await _basic_layout_adjustment_async(presentation_path)

        # 在线程池中执行COM操作以避免阻塞
        loop = asyncio.get_event_loop()

        def _call_set_textbox_autofit():
            """调用现有的set_textbox_autofit函数"""
            try:
                # 导入set_textbox_autofit函数
                from .adjust_text_size import set_textbox_autofit

                # 获取绝对路径
                abs_path = os.path.abspath(presentation_path)
                logger.debug(f"调用set_textbox_autofit，文件路径: {abs_path}")

                # 调用现有的布局调整函数
                result = set_textbox_autofit(abs_path)

                if result:
                    logger.info("set_textbox_autofit调用成功")
                    return True
                else:
                    logger.warning("set_textbox_autofit调用失败")
                    return False

            except ImportError as import_error:
                logger.warning(f"无法导入set_textbox_autofit函数: {import_error}")
                return False
            except Exception as e:
                logger.error(f"调用set_textbox_autofit时出错: {e}")
                return False

        # 在线程池中执行COM操作
        result = await loop.run_in_executor(None, _call_set_textbox_autofit)

        # 如果COM操作失败，尝试基础调整
        if not result:
            logger.info("set_textbox_autofit调用失败，尝试基础布局调整")
            return await _basic_layout_adjustment_async(presentation_path)

        return result

    except Exception as e:
        logger.error(f"布局调整过程出错: {e}")
        # 尝试基础调整作为后备方案
        return await _basic_layout_adjustment_async(presentation_path)

async def _basic_layout_adjustment_async(presentation_path: str) -> bool:
    """
    基础布局调整（不依赖COM，跨平台兼容）
    确保所有文本框都设置为MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

    Args:
        presentation_path: PPT文件路径

    Returns:
        调整是否成功
    """
    try:
        loop = asyncio.get_event_loop()

        def _basic_adjustment():
            try:
                # 加载演示文稿
                prs = Presentation(presentation_path)

                total_textboxes = 0
                processed_textboxes = 0

                # 遍历所有幻灯片
                for slide_index, slide in enumerate(prs.slides, 1):
                    logger.debug(f"处理第 {slide_index} 张幻灯片的文本框自动调整")

                    for shape in slide.shapes:
                        # 处理普通文本框
                        if shape.has_text_frame:
                            total_textboxes += 1
                            try:
                                # 使用内容检测的增强自适应设置
                                text_frame = shape.text_frame
                                result = safe_set_autofit_with_content_check(text_frame, shape, debug=True)

                                if result['adjusted']:
                                    processed_textboxes += 1
                                    logger.debug(f"已处理文本框: 幻灯片{slide_index}, 内容: {result['content']}")
                                else:
                                    logger.debug(f"跳过文本框: 幻灯片{slide_index}, 原因: {result['reason']}")

                            except Exception as shape_error:
                                logger.warning(f"处理幻灯片{slide_index}的文本框时出错: {shape_error}")

                        # 处理表格中的文本框
                        elif shape.has_table:
                            table = shape.table
                            logger.debug(f"处理第 {slide_index} 张幻灯片的表格 ({table.rows} 行 x {table.columns} 列)")

                            for row_index, row in enumerate(table.rows):
                                for col_index, cell in enumerate(row.cells):
                                    total_textboxes += 1
                                    try:
                                        # 只设置表格单元格文本框自动调整，不改变其他格式
                                        text_frame = cell.text_frame
                                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                                        processed_textboxes += 1
                                        logger.debug(f"已处理表格单元格: 幻灯片{slide_index}, 行{row_index+1}, 列{col_index+1}")

                                    except Exception as cell_error:
                                        logger.warning(f"处理幻灯片{slide_index}表格单元格({row_index+1},{col_index+1})时出错: {cell_error}")

                # 保存演示文稿
                prs.save(presentation_path)
                logger.info(f"基础布局调整完成: 处理了 {processed_textboxes}/{total_textboxes} 个文本框")
                return True

            except Exception as e:
                logger.error(f"基础布局调整失败: {e}")
                return False

        # 在线程池中执行文件操作
        return await loop.run_in_executor(None, _basic_adjustment)

    except Exception as e:
        logger.error(f"基础布局调整过程出错: {e}")
        return False

async def ensure_all_textboxes_autofit_async(presentation_path: str) -> bool:
    """
    确保PPT中所有文本框都设置为自动调整大小
    这是一个专门的函数，用于解决文本框未全部设置为MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE的问题

    Args:
        presentation_path: PPT文件路径

    Returns:
        调整是否成功
    """
    try:
        loop = asyncio.get_event_loop()

        def _ensure_all_autofit():
            try:
                # 加载演示文稿
                prs = Presentation(presentation_path)

                total_shapes = 0
                total_textboxes = 0
                processed_textboxes = 0
                skipped_shapes = 0

                logger.info("开始确保所有文本框都设置为自动调整大小...")

                # 遍历所有幻灯片
                for slide_index, slide in enumerate(prs.slides, 1):
                    logger.debug(f"检查第 {slide_index} 张幻灯片的所有形状...")

                    for shape_index, shape in enumerate(slide.shapes):
                        total_shapes += 1

                        try:
                            # 处理普通文本框
                            if shape.has_text_frame:
                                total_textboxes += 1
                                text_frame = shape.text_frame

                                # 使用内容检测的增强自适应设置
                                result = safe_set_autofit_with_content_check(text_frame, shape, debug=True)
                                text_frame.word_wrap = True

                                if result['adjusted']:
                                    processed_textboxes += 1
                                    logger.debug(f"  幻灯片{slide_index}-形状{shape_index+1}: 已设置文本框自动调整，内容: {result['content']}")
                                else:
                                    logger.debug(f"跳过幻灯片{slide_index}-形状{shape_index+1}: {result['reason']}")

                            # 处理表格
                            elif shape.has_table:
                                table = shape.table
                                logger.debug(f"处理表格: {table.rows} 行 x {table.columns} 列")

                                for row_index, row in enumerate(table.rows):
                                    for col_index, cell in enumerate(row.cells):
                                        total_textboxes += 1

                                        # 表格单元格的文本框
                                        text_frame = cell.text_frame
                                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                                        text_frame.word_wrap = True

                                        processed_textboxes += 1
                                        logger.debug(f"  幻灯片{slide_index}-表格单元格({row_index+1},{col_index+1}): 已设置自动调整")

                            else:
                                skipped_shapes += 1
                                logger.debug(f"跳过非文本形状: 幻灯片{slide_index}-形状{shape_index+1} (类型: {shape.shape_type})")

                        except Exception as shape_error:
                            logger.warning(f"处理幻灯片{slide_index}-形状{shape_index+1}时出错: {shape_error}")

                # 保存演示文稿
                prs.save(presentation_path)

                logger.info(f"文本框自动调整设置完成:")
                logger.info(f"  - 总形状数: {total_shapes}")
                logger.info(f"  - 文本框总数: {total_textboxes}")
                logger.info(f"  - 已处理文本框: {processed_textboxes}")
                logger.info(f"  - 跳过的形状: {skipped_shapes}")
                logger.info(f"  - 成功率: {(processed_textboxes/total_textboxes*100):.1f}%" if total_textboxes > 0 else "  - 成功率: N/A")

                return True

            except Exception as e:
                logger.error(f"确保文本框自动调整失败: {e}")
                import traceback
                logger.error(f"错误详情: {traceback.format_exc()}")
                return False

        # 在线程池中执行文件操作
        return await loop.run_in_executor(None, _ensure_all_autofit)

    except Exception as e:
        logger.error(f"确保文本框自动调整过程出错: {e}")
        return False

async def _preserve_textbox_size_with_autofit_async(presentation_path: str) -> bool:
    """
    异步设置文本框自适应并保持原始大小

    Args:
        presentation_path: PPT文件路径

    Returns:
        调整是否成功
    """
    try:
        loop = asyncio.get_event_loop()

        def _preserve_size_autofit():
            try:
                # 加载演示文稿
                prs = Presentation(presentation_path)

                total_textboxes = 0
                processed_textboxes = 0
                size_preserved_count = 0

                logger.info("开始设置文本框自适应并保持原始大小...")

                # 遍历所有幻灯片
                for slide_index, slide in enumerate(prs.slides, 1):
                    logger.debug(f"处理第 {slide_index} 张幻灯片的文本框...")

                    for shape_index, shape in enumerate(slide.shapes):
                        try:
                            # 处理普通文本框
                            if shape.has_text_frame:
                                total_textboxes += 1

                                # 记录原始尺寸
                                original_width = shape.width
                                original_height = shape.height
                                original_left = shape.left
                                original_top = shape.top

                                # 使用内容检测的增强自适应设置
                                text_frame = shape.text_frame
                                result = safe_set_autofit_with_content_check(text_frame, shape, debug=True)

                                if not result['adjusted']:
                                    logger.debug(f"跳过文本框: 幻灯片{slide_index}-形状{shape_index+1}, 原因: {result['reason']}")
                                    continue

                                # 检查并恢复原始尺寸
                                if (shape.width != original_width or
                                    shape.height != original_height or
                                    shape.left != original_left or
                                    shape.top != original_top):

                                    shape.width = original_width
                                    shape.height = original_height
                                    shape.left = original_left
                                    shape.top = original_top

                                    size_preserved_count += 1
                                    logger.debug(f"  已恢复文本框原始尺寸: 幻灯片{slide_index}-形状{shape_index+1}")

                                processed_textboxes += 1
                                logger.debug(f"  幻灯片{slide_index}-形状{shape_index+1}: 已设置文本框自适应")

                            # 处理表格
                            elif shape.has_table:
                                table = shape.table
                                logger.debug(f"处理表格: {table.rows} 行 x {table.columns} 列")

                                # 记录表格原始尺寸
                                table_original_width = shape.width
                                table_original_height = shape.height
                                table_original_left = shape.left
                                table_original_top = shape.top

                                for row_index, row in enumerate(table.rows):
                                    for col_index, cell in enumerate(row.cells):
                                        total_textboxes += 1

                                        # 只设置表格单元格文本框自适应，不改变其他格式
                                        text_frame = cell.text_frame
                                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                                        processed_textboxes += 1
                                        logger.debug(f"  幻灯片{slide_index}-表格单元格({row_index+1},{col_index+1}): 已设置自适应")

                                # 确保表格整体尺寸不变
                                if (shape.width != table_original_width or
                                    shape.height != table_original_height or
                                    shape.left != table_original_left or
                                    shape.top != table_original_top):

                                    shape.width = table_original_width
                                    shape.height = table_original_height
                                    shape.left = table_original_left
                                    shape.top = table_original_top

                                    size_preserved_count += 1
                                    logger.debug(f"  已恢复表格原始尺寸: 幻灯片{slide_index}-表格")

                        except Exception as shape_error:
                            logger.warning(f"处理幻灯片{slide_index}-形状{shape_index+1}时出错: {shape_error}")

                # 保存演示文稿
                prs.save(presentation_path)

                logger.info(f"文本框自适应设置完成（保持原始大小）:")
                logger.info(f"  - 文本框总数: {total_textboxes}")
                logger.info(f"  - 已处理文本框: {processed_textboxes}")
                logger.info(f"  - 尺寸保护次数: {size_preserved_count}")
                logger.info(f"  - 成功率: {(processed_textboxes/total_textboxes*100):.1f}%" if total_textboxes > 0 else "  - 成功率: N/A")

                return True

            except Exception as e:
                logger.error(f"设置文本框自适应失败: {e}")
                return False

        # 在线程池中执行文件操作
        return await loop.run_in_executor(None, _preserve_size_autofit)

    except Exception as e:
        logger.error(f"设置文本框自适应过程出错: {e}")
        return False











async def _unified_shape_processing_async(presentation_path: str) -> bool:
    """
    统一的形状处理函数（避免多重处理冲突）
    集成布局调整、自适应设置、尺寸保护等功能

    Args:
        presentation_path: PPT文件路径

    Returns:
        处理是否成功
    """
    try:
        loop = asyncio.get_event_loop()

        def _unified_processing():
            try:
                # 加载演示文稿
                prs = Presentation(presentation_path)

                total_shapes = 0
                total_textboxes = 0
                processed_textboxes = 0
                skipped_textboxes = 0
                protected_shapes = 0

                logger.info("开始统一形状处理...")

                # 遍历所有幻灯片
                for slide_index, slide in enumerate(prs.slides, 1):
                    logger.debug(f"处理第 {slide_index} 张幻灯片的形状...")

                    for shape_index, shape in enumerate(slide.shapes):
                        total_shapes += 1

                        try:
                            # 处理普通文本框
                            if shape.has_text_frame:
                                total_textboxes += 1
                                text_frame = shape.text_frame

                                # 使用增强的内容检测和形状保护
                                result = safe_set_autofit_with_content_check(text_frame, shape, debug=True)

                                if result['adjusted']:
                                    if result.get('success', True):
                                        processed_textboxes += 1
                                        logger.debug(f"  幻灯片{slide_index}-形状{shape_index+1}: 处理成功")
                                    else:
                                        protected_shapes += 1
                                        logger.debug(f"🛡️ 幻灯片{slide_index}-形状{shape_index+1}: 检测到变形，已保护")
                                else:
                                    skipped_textboxes += 1
                                    logger.debug(f"⏭️ 幻灯片{slide_index}-形状{shape_index+1}: 跳过处理")

                            # 处理表格（表格单元格使用标准处理）
                            elif shape.has_table:
                                table = shape.table
                                for row_index, row in enumerate(table.rows):
                                    for col_index, cell in enumerate(row.cells):
                                        total_textboxes += 1

                                        # 表格单元格使用标准自适应（性能考虑）
                                        text_frame = cell.text_frame
                                        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

                                        processed_textboxes += 1
                                        logger.debug(f"  幻灯片{slide_index}-表格单元格({row_index+1},{col_index+1}): 已设置自适应")

                        except Exception as shape_error:
                            logger.warning(f"处理幻灯片{slide_index}-形状{shape_index+1}时出错: {shape_error}")

                # 保存演示文稿
                prs.save(presentation_path)

                logger.info(f"统一形状处理完成:")
                logger.info(f"  - 形状总数: {total_shapes}")
                logger.info(f"  - 文本框总数: {total_textboxes}")
                logger.info(f"  - 已处理文本框: {processed_textboxes}")
                logger.info(f"  - 跳过文本框: {skipped_textboxes}")
                logger.info(f"  - 保护形状: {protected_shapes}")
                logger.info(f"  - 处理成功率: {(processed_textboxes/max(total_textboxes, 1)*100):.1f}%")

                return True

            except Exception as e:
                logger.error(f"统一形状处理失败: {e}")
                return False

        # 在线程池中执行文件操作
        return await loop.run_in_executor(None, _unified_processing)

    except Exception as e:
        logger.error(f"统一形状处理过程出错: {e}")
        return False


async def process_presentation_async(presentation_path: str,
                                   stop_words_list: List[str],
                                   custom_translations: Dict[str, str],
                                   select_page: List[int],
                                   source_language: str,
                                   target_language: str,
                                   bilingual_translation: str,
                                   progress_callback=None) -> bool:
    """
    异步处理演示文稿（基于页面的翻译机制）
    每页调用一次API，按段落匹配翻译结果

    Args:
        presentation_path: PPT文件路径
        stop_words_list: 停止词列表
        custom_translations: 自定义翻译字典
        select_page: 选中的页面列表
        source_language: 源语言代码
        target_language: 目标语言代码
        bilingual_translation: 是否双语翻译
        progress_callback: 进度回调函数，接收两个参数(current_slide, total_slides)

    Returns:
        处理是否成功
    """
    start_time = time.time()
    logger.info(f"开始异步处理演示文稿: {os.path.basename(presentation_path)}")
    logger.info(f"源语言: {source_language}, 目标语言: {target_language}, 双语翻译: {bilingual_translation}")
    logger.info(f"选中页面: {select_page}")

    try:
        # 加载演示文稿
        logger.info("正在加载演示文稿...")
        loop = asyncio.get_event_loop()

        def _read_presentation():
            return Presentation(presentation_path)

        prs = await loop.run_in_executor(None, _read_presentation)
        total_slides = len(prs.slides)
        logger.info(f"演示文稿加载成功，共 {total_slides} 张幻灯片")
        logger.info(f" 选择的页面参数: {select_page}")

        # 如果没有选择页面，默认翻译所有页面
        if not select_page:
            select_page = list(range(1, total_slides + 1))
            logger.info(f" 没有指定页面，将翻译所有页面: {select_page}")
        else:
            logger.info(f" 将翻译指定页面: {select_page}")

        # 获取领域（使用第一页的内容分析）
        logger.info("正在分析文本领域...")
        first_slide_text = ""
        if prs.slides:
            for shape in prs.slides[0].shapes:
                if shape.has_text_frame:
                    first_slide_text += shape.text_frame.text + "\n"

        field = await get_field_async(first_slide_text[:500])  # 只用前500字符分析领域
        logger.info(f"文本领域分析结果: {field}")

        # 初始化进度
        if progress_callback:
            progress_callback(0, total_slides)

        # 处理每张幻灯片
        processed_slides = 0
        skipped_slides = 0
        total_translated_paragraphs = 0

        for current_slide_index, slide in enumerate(prs.slides, 1):
            # 更新翻译进度
            if progress_callback:
                progress_callback(current_slide_index - 1, total_slides)

            # 检查是否需要处理当前幻灯片
            if select_page and current_slide_index not in select_page:
                logger.info(f"跳过第 {current_slide_index} 张幻灯片 (不在选中页面列表中)")
                skipped_slides += 1
                continue

            # 使用基于页面的翻译
            logger.info(f"开始处理第 {current_slide_index} 张幻灯片...")
            slide_start_time = time.time()

            translated_count = await translate_slide_by_page(
                slide, current_slide_index - 1, source_language, target_language,
                bilingual_translation, field
            )

            slide_elapsed = time.time() - slide_start_time
            total_translated_paragraphs += translated_count
            processed_slides += 1

            logger.info(f"第 {current_slide_index} 张幻灯片处理完成，翻译了 {translated_count} 个段落，耗时: {slide_elapsed:.2f}秒")

        # 保存演示文稿
        logger.info("正在保存演示文稿...")

        def _save_presentation():
            # 创建临时文件以避免内存泄漏
            temp_path = f"{presentation_path}.temp"
            # 使用已经修改过的演示文稿对象进行保存
            prs.save(temp_path)

            # 如果保存成功，替换原文件
            if os.path.exists(temp_path):
                if os.path.exists(presentation_path):
                    os.remove(presentation_path)
                os.rename(temp_path, presentation_path)
                return True
            return False

        save_result = await loop.run_in_executor(None, _save_presentation)

        # 如果保存成功，进行简单的布局调整
        if save_result:
            logger.info("正在进行布局调整...")

            # 使用COM操作进行最终的文本框调整
            layout_result = await _adjust_ppt_layout_async(presentation_path)
            if layout_result:
                logger.info("布局调整完成")
            else:
                logger.warning("布局调整失败，但翻译已完成")

        elapsed = time.time() - start_time
        logger.info(f"演示文稿处理完成:")
        logger.info(f"  - 处理了 {processed_slides} 张幻灯片")
        logger.info(f"  - 跳过了 {skipped_slides} 张幻灯片")
        logger.info(f"  - 翻译了 {total_translated_paragraphs} 个段落")
        logger.info(f"  - 总耗时: {elapsed:.2f}秒")
        logger.info(f"  - 平均每页耗时: {elapsed/max(processed_slides, 1):.2f}秒")

        # 最终更新进度为100%
        if progress_callback:
            progress_callback(total_slides, total_slides)

        return save_result
    except Exception as e:
        logger.error(f"处理演示文稿时出错: {str(e)}")
        import traceback
        logger.error(traceback.format_exc())

        # 在出错时也更新进度
        if progress_callback:
            progress_callback(0, 1)

        return False

async def process_presentation_add_annotations_async(presentation_path: str,
                                                 annotations: Dict,
                                                 stop_words: List[str],
                                                 custom_translations: Dict[str, str],
                                                 source_language: str,
                                                 target_language: str,
                                                 bilingual_translation: str,
                                                 progress_callback=None) -> bool:
    """
    异步处理带注释的演示文稿

    Args:
        presentation_path: PPT文件路径
        annotations: 注释数据
        stop_words: 停止词列表
        custom_translations: 自定义翻译字典
        source_language: 源语言代码
        target_language: 目标语言代码
        bilingual_translation: 是否双语翻译
        progress_callback: 进度回调函数，接收两个参数(current_slide, total_slides)

    Returns:
        处理是否成功
    """
    start_time = time.time()
    logger.info(f"开始异步处理带注释的演示文稿: {presentation_path}")

    try:
        # 检查注释数据
        if not annotations or "annotations" not in annotations:
            logger.error("无效的注释数据")
            return False

        annotation_items = annotations.get("annotations", [])
        total_annotations = len(annotation_items)
        logger.info(f"共有 {total_annotations} 个注释项")

        # 初始化进度
        if progress_callback:
            progress_callback(0, total_annotations)

        # 使用线程池执行文件IO操作
        loop = asyncio.get_event_loop()

        async def _apply_annotations():
            try:
                # 加载演示文稿
                prs = Presentation(presentation_path)

                # 检查幻灯片数量
                if len(prs.slides) == 0:
                    logger.error("演示文稿中没有幻灯片")
                    return False

                # 收集所有注释文本进行翻译
                logger.info("正在收集注释文本...")
                all_text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text_frame = shape.text_frame
                            for paragraph in text_frame.paragraphs:
                                text = paragraph.text.strip()
                                if text:
                                    all_text += text + "\n"

                # 获取领域（使用新的阿里云异步API）
                from .local_qwen_async import get_field_async
                field = await get_field_async(all_text)
                logger.info(f"文本领域分析结果: {field}")

                # 准备注释文本进行翻译
                tage_text = ""
                annotation_items = annotations["annotations"]
                for item in annotation_items:
                    text = item["ocrResult"].replace("\n", " ")
                    tage_text += text + "\n"

                # 处理停止词和自定义翻译
                stop_words_filtered = []
                custom_words = {}
                for word in stop_words:
                    if word in tage_text:
                        stop_words_filtered.append(word)
                for k, v in custom_translations.items():
                    if k in tage_text:
                        custom_words[k] = v

                # 翻译注释文本（使用新的阿里云异步API）
                logger.info("正在翻译注释文本...")
                from .local_qwen_async import translate_async
                from .ppt_translate import find_most_similar
                data = await translate_async(tage_text, field, stop_words_filtered, custom_words, source_language, target_language)
                logger.info(f"翻译完成，共翻译 {len(data)} 个文本段")

                # 处理每个注释，添加到对应页面右上角
                processed_count = 0
                for i, item in enumerate(annotation_items):
                    try:
                        page = item["page"]
                        original_text = item["ocrResult"]

                        # 页面索引从1开始，转换为0开始
                        slide_index = page - 1
                        if slide_index < 0 or slide_index >= len(prs.slides):
                            logger.warning(f"页面索引超出范围: {page}, 跳过此注释")
                            continue

                        slide = prs.slides[slide_index]
                        slide_width = prs.slide_width
                        slide_height = prs.slide_height

                        # 在右上角添加翻译文本框
                        left = slide_width - Inches(2)  # 右边距2英寸
                        top = 0  # 顶部
                        width = Inches(2)  # 宽度2英寸
                        height = Inches(1)  # 高度1英寸

                        # 查找最相似的翻译
                        new_text = find_most_similar(original_text, list(data.keys()))
                        if new_text in data:
                            translated_text = data[new_text]

                            # 添加文本框
                            textbox = slide.shapes.add_textbox(left, top, width, height)
                            text_frame = textbox.text_frame
                            text_frame.text = translated_text

                            # 设置字体为红色
                            for paragraph in text_frame.paragraphs:
                                for run in paragraph.runs:
                                    run.font.color.rgb = RGBColor(255, 0, 0)  # 红色
                                    run.font.size = Pt(12)  # 设置字体大小

                            # 设置文本框自适应
                            text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                            text_frame.word_wrap = True

                            processed_count += 1
                            logger.info(f"  第 {page} 页添加翻译注释: '{original_text[:30]}...' -> '{translated_text[:30]}...'")
                        else:
                            logger.warning(f"未找到匹配的翻译: {original_text[:30]}...")

                        # 更新进度
                        if progress_callback:
                            progress_callback(i + 1, total_annotations)

                    except Exception as e:
                        logger.error(f"处理注释 {i} 时出错: {str(e)}")
                        continue

                # 保存演示文稿
                temp_path = f"{presentation_path}.temp"
                prs.save(temp_path)

                # 如果保存成功，替换原文件
                if os.path.exists(temp_path):
                    if os.path.exists(presentation_path):
                        os.remove(presentation_path)
                    os.rename(temp_path, presentation_path)

                logger.info(f"处理了 {processed_count}/{total_annotations} 个注释")
                return True

            except Exception as e:
                logger.error(f"应用注释时出错: {str(e)}")
                import traceback
                logger.error(f"错误详情: {traceback.format_exc()}")
                return False

        # 异步应用注释
        result = await _apply_annotations()

        elapsed = time.time() - start_time
        logger.info(f"带注释的演示文稿处理完成，耗时: {elapsed:.2f}秒")

        # 最终更新进度为100%
        if progress_callback:
            progress_callback(total_annotations, total_annotations)

        return result
    except Exception as e:
        logger.error(f"处理带注释的演示文稿失败: {str(e)}")
        import traceback
        logger.error(f"错误详情: {traceback.format_exc()}")

        # 在出错时也更新进度
        if progress_callback:
            progress_callback(0, 1)

        raise

# 创建同步包装函数
def run_async_in_thread(func, *args, **kwargs):
    """
    在线程中运行异步函数

    Args:
        func: 异步函数
        *args: 位置参数
        **kwargs: 关键字参数

    Returns:
        函数结果
    """
    import asyncio
    import threading
    import concurrent.futures

    # 检查是否已有运行中的事件循环
    try:
        current_loop = asyncio.get_running_loop()
        # 如果有运行中的循环，在新线程中运行
        def run_in_new_thread():
            new_loop = asyncio.new_event_loop()
            asyncio.set_event_loop(new_loop)
            try:
                return new_loop.run_until_complete(func(*args, **kwargs))
            finally:
                new_loop.close()

        with concurrent.futures.ThreadPoolExecutor() as executor:
            future = executor.submit(run_in_new_thread)
            return future.result()
    except RuntimeError:
        # 没有运行中的循环，直接运行
        try:
            loop = asyncio.get_event_loop()
            if loop.is_closed():
                raise RuntimeError("Event loop is closed")
        except RuntimeError:
            loop = asyncio.new_event_loop()
            asyncio.set_event_loop(loop)

        try:
            return loop.run_until_complete(func(*args, **kwargs))
        finally:
            # 不关闭循环，让其他代码继续使用
            pass

# 创建同步API（兼容原始接口）
def process_presentation(presentation_path: str,
                       stop_words_list: List[str] = None,
                       custom_translations: Dict[str, str] = None,
                       select_page: List[int] = None,
                       source_language: str = "英语",
                       target_language: str = "中文",
                       bilingual_translation: str = "0",
                       progress_callback=None,
                       # 兼容性参数
                       stop_words: List[str] = None,
                       **kwargs) -> bool:
    """
    处理PPT翻译（同步包装函数）
    使用基于页面的翻译机制，每页调用一次API，按段落匹配翻译结果

    Args:
        presentation_path: PPT文件路径
        stop_words_list: 停止词列表（优先使用）
        custom_translations: 自定义翻译字典
        select_page: 选中的页面列表
        source_language: 源语言代码
        target_language: 目标语言代码
        bilingual_translation: 是否双语翻译（"0"或"1"）
        progress_callback: 进度回调函数，接收两个参数(current_slide, total_slides)
        stop_words: 停止词列表（兼容性参数）

    Returns:
        处理是否成功
    """
    logger.info(f"开始处理演示文稿: {os.path.basename(presentation_path)}")
    logger.info(f"参数: 源语言={source_language}, 目标语言={target_language}, 双语={bilingual_translation}")

    # 处理兼容性参数
    if stop_words_list is None and stop_words is not None:
        stop_words_list = stop_words
    elif stop_words_list is None:
        stop_words_list = []

    if custom_translations is None:
        custom_translations = {}

    if select_page is None:
        select_page = []

    try:
        # 使用基于页面的翻译模式
        result = run_async_in_thread(
            process_presentation_async,
            presentation_path,
            stop_words_list,
            custom_translations,
            select_page,
            source_language,
            target_language,
            bilingual_translation,
            progress_callback
        )

        logger.info(f"演示文稿处理完成: {os.path.basename(presentation_path)}")
        return result
    except Exception as e:
        logger.error(f"处理演示文稿失败: {os.path.basename(presentation_path)}, 错误: {str(e)}")
        import traceback
        logger.error(f"错误详情: {traceback.format_exc()}")
        return False

def process_presentation_add_annotations(presentation_path: str,
                                       annotations: Dict,
                                       stop_words: List[str],
                                       custom_translations: Dict[str, str],
                                       source_language: str,
                                       target_language: str,
                                       bilingual_translation: str,
                                       progress_callback=None) -> bool:
    """
    处理带注释的PPT翻译（同步包装函数）

    Args:
        presentation_path: PPT文件路径
        annotations: 注释字典
        stop_words: 停止词列表
        custom_translations: 自定义翻译字典
        source_language: 源语言代码
        target_language: 目标语言代码
        bilingual_translation: 是否双语翻译（"0"或"1"）
        progress_callback: 进度回调函数，接收两个参数(current_slide, total_slides)

    Returns:
        处理是否成功
    """
    logger.info(f"开始处理带注释的演示文稿: {presentation_path}")

    # 转换双语翻译参数
    is_bilingual = bilingual_translation == "1"

    try:
        # 在线程中运行异步函数
        result = run_async_in_thread(
            process_presentation_add_annotations_async,
            presentation_path,
            annotations,
            stop_words,
            custom_translations,
            source_language,
            target_language,
            is_bilingual,
            progress_callback
        )
        logger.info(f"带注释的演示文稿处理完成: {presentation_path}")
        return result
    except Exception as e:
        logger.error(f"处理带注释的演示文稿失败: {presentation_path}, 错误: {str(e)}")
        import traceback
        logger.error(f"错误详情: {traceback.format_exc()}")
        return False